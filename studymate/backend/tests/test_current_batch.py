from __future__ import annotations

import tempfile
import unittest
from io import BytesIO
from pathlib import Path

from fastapi import HTTPException
from sqlalchemy.ext.asyncio import async_sessionmaker, create_async_engine

from app.api.knowledge import (
    KnowledgeBaseCreate,
    _extract_pages,
    _split_pages,
    create_library,
    list_libraries,
    search_owned_library,
)
from app.api.tutor import TutorChatRequest, TutorMessage, tutor_chat, tutor_models
from app.db import models  # noqa: F401 - register all tables
from app.db.models import User, UserKnowledgeChunk, UserKnowledgeDocument
from app.db.session import Base


class CurrentBatchApiTests(unittest.IsolatedAsyncioTestCase):
    async def asyncSetUp(self):
        self._temp_dir = tempfile.TemporaryDirectory(prefix="studymate-current-batch-")
        db_path = Path(self._temp_dir.name) / "test.db"
        self._engine = create_async_engine(f"sqlite+aiosqlite:///{db_path}")
        self._sessions = async_sessionmaker(self._engine, expire_on_commit=False)
        async with self._engine.begin() as connection:
            await connection.run_sync(Base.metadata.create_all)

    async def asyncTearDown(self):
        await self._engine.dispose()
        self._temp_dir.cleanup()

    async def test_private_library_search_never_crosses_users(self):
        async with self._sessions() as db:
            first = User(name="first", email="first@example.test")
            second = User(name="second", email="second@example.test")
            db.add_all([first, second])
            await db.commit()
            await db.refresh(first)
            await db.refresh(second)

            library = await create_library(KnowledgeBaseCreate(name="私有教材"), first, db)
            document = UserKnowledgeDocument(
                user_id=first.id,
                knowledge_base_id=library["id"],
                filename="chapter.md",
                media_type="text/markdown",
                size=32,
                status="ready_keyword",
                parse_progress=100,
                vector_progress=0,
                page_count=1,
            )
            db.add(document)
            await db.flush()
            db.add(UserKnowledgeChunk(
                user_id=first.id,
                knowledge_base_id=library["id"],
                document_id=document.id,
                content="梯度下降沿损失函数的负梯度方向更新参数。",
                page=3,
            ))
            await db.commit()

            own_hits = await search_owned_library(
                db,
                user_id=first.id,
                library_id=library["id"],
                query="梯度下降",
                with_semantic=False,
            )
            other_hits = await search_owned_library(
                db,
                user_id=second.id,
                library_id=library["id"],
                query="梯度下降",
                with_semantic=False,
            )
            first_list = await list_libraries(first, db)
            second_list = await list_libraries(second, db)

            self.assertEqual(own_hits[0]["source"], "chapter.md")
            self.assertEqual(own_hits[0]["page"], 3)
            self.assertEqual(other_hits, [])
            self.assertEqual(first_list["count"], 1)
            self.assertEqual(second_list["count"], 0)

    async def test_model_catalog_exposes_status_but_not_credentials(self):
        payload = await tutor_models()
        self.assertEqual([item["id"] for item in payload["items"]], ["qwen", "deepseek", "spark", "mimo"])
        serialized = repr(payload).lower()
        self.assertNotIn("api_key", serialized)
        self.assertNotIn("base_url", serialized)

        with self.assertRaises(HTTPException) as raised:
            await tutor_chat(
                TutorChatRequest(
                    provider="deepseek",
                    messages=[
                        TutorMessage(
                            role="user",
                            content="解释图片",
                            images=["data:image/png;base64,AA=="],
                        )
                    ],
                ),
                User(id=1, name="student"),
            )
        self.assertEqual(raised.exception.status_code, 422)


class CurrentBatchParserTests(unittest.TestCase):
    def test_text_chunks_preserve_page_and_size_limit(self):
        chunks = _split_pages([(7, "第一段知识。\n" + "长" * 25)], max_chars=12)
        self.assertTrue(chunks)
        self.assertTrue(all(page == 7 for page, _ in chunks))
        self.assertTrue(all(len(content) <= 12 for _, content in chunks))

    def test_pptx_parser_preserves_slide_numbers(self):
        from pptx import Presentation

        output = BytesIO()
        presentation = Presentation()
        first = presentation.slides.add_slide(presentation.slide_layouts[1])
        first.shapes.title.text = "第一页标题"
        first.placeholders[1].text = "第一页内容"
        second = presentation.slides.add_slide(presentation.slide_layouts[1])
        second.shapes.title.text = "第二页标题"
        second.placeholders[1].text = "第二页内容"
        presentation.save(output)

        pages = _extract_pages("lesson.pptx", output.getvalue())
        self.assertEqual([page for page, _ in pages], [1, 2])
        self.assertIn("第一页内容", pages[0][1])
        self.assertIn("第二页内容", pages[1][1])


if __name__ == "__main__":
    unittest.main()
