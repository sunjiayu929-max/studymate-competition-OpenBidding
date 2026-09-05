"""用户私有知识库：创建、资料解析、向量化进度、检索与来源追溯。"""
from __future__ import annotations

import hashlib
from io import BytesIO
from datetime import datetime
from pathlib import Path

from fastapi import APIRouter, BackgroundTasks, Depends, File, HTTPException, Query, UploadFile, status
from pydantic import BaseModel, Field
from sqlalchemy import delete, select
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.config import settings
from app.db import get_db
from app.db.session import async_session_maker
from app.db.models import (
    User,
    UserKnowledgeBase,
    UserKnowledgeChunk,
    UserKnowledgeDocument,
)
from app.deps import require_user
from app.llm.embeddings import embed_query, embed_texts, has_embedding_key
from app.rag.engine import BM25Engine, Chunk, VectorIndex

router = APIRouter(prefix="/knowledge-bases", tags=["knowledge-bases"])

MAX_FILE_BYTES = 20 * 1024 * 1024
SUPPORTED_SUFFIXES = {".pdf", ".ppt", ".pptx", ".doc", ".docx", ".md", ".markdown", ".txt"}
MAX_SAFE_RETRIES = 3


class KnowledgeBaseCreate(BaseModel):
    name: str = Field(min_length=1, max_length=128)
    description: str = Field(default="", max_length=800)
    bound_course_id: int | None = None


class KnowledgeBaseUpdate(BaseModel):
    name: str | None = Field(default=None, min_length=1, max_length=128)
    description: str | None = Field(default=None, max_length=800)
    bound_course_id: int | None = None


def _library_payload(row: UserKnowledgeBase, documents: list[UserKnowledgeDocument] | None = None) -> dict:
    docs = documents or []
    return {
        "id": row.id,
        "name": row.name,
        "description": row.description,
        "bound_course_id": row.bound_course_id,
        "created_at": row.created_at.isoformat() if row.created_at else None,
        "updated_at": row.updated_at.isoformat() if row.updated_at else None,
        "document_count": len(docs),
        "documents": [
            {
                "id": doc.id,
                "filename": doc.filename,
                "media_type": doc.media_type,
                "size": doc.size,
                "status": doc.status,
                "parse_progress": doc.parse_progress,
                "vector_progress": doc.vector_progress,
                "error_detail": doc.error_detail,
                "page_count": doc.page_count,
                "retry_count": doc.retry_count,
                "retry_available": doc.status == "error" and doc.retry_count < MAX_SAFE_RETRIES and bool(doc.source_path),
                "ocr_status": doc.ocr_status,
                "started_at": doc.started_at.isoformat() if doc.started_at else None,
                "finished_at": doc.finished_at.isoformat() if doc.finished_at else None,
                "created_at": doc.created_at.isoformat() if doc.created_at else None,
            }
            for doc in docs
        ],
    }


async def _owned_library(db: AsyncSession, library_id: int, user_id: int) -> UserKnowledgeBase:
    row = await db.scalar(
        select(UserKnowledgeBase).where(
            UserKnowledgeBase.id == library_id,
            UserKnowledgeBase.user_id == user_id,
        )
    )
    if row is None:
        raise HTTPException(status_code=404, detail="知识库不存在")
    return row


def _storage_root() -> Path:
    root = Path(settings.PRIVATE_KNOWLEDGE_DIR).expanduser().resolve()
    root.mkdir(parents=True, exist_ok=True)
    return root


def _source_path(*, user_id: int, document_id: int, filename: str) -> Path:
    suffix = Path(filename).suffix.lower()
    directory = (_storage_root() / str(user_id)).resolve()
    directory.mkdir(parents=True, exist_ok=True)
    return directory / f"{document_id}{suffix}"


def _safe_unlink(path_text: str) -> None:
    if not path_text:
        return
    try:
        path = Path(path_text).resolve()
        root = _storage_root()
        if path.is_relative_to(root) and path.is_file():
            path.unlink()
    except (OSError, ValueError):
        # 数据库清理不能因本地文件已经丢失而失败。
        return


@router.get("")
async def list_libraries(user: User = Depends(require_user), db: AsyncSession = Depends(get_db)):
    rows = (
        await db.scalars(
            select(UserKnowledgeBase)
            .where(UserKnowledgeBase.user_id == user.id)
            .order_by(UserKnowledgeBase.updated_at.desc(), UserKnowledgeBase.id.desc())
        )
    ).all()
    documents = (
        await db.scalars(
            select(UserKnowledgeDocument)
            .where(UserKnowledgeDocument.user_id == user.id)
            .order_by(UserKnowledgeDocument.id.desc())
        )
    ).all()
    grouped: dict[int, list[UserKnowledgeDocument]] = {}
    for document in documents:
        grouped.setdefault(document.knowledge_base_id, []).append(document)
    return {"count": len(rows), "items": [_library_payload(row, grouped.get(row.id, [])) for row in rows]}


@router.post("")
async def create_library(
    payload: KnowledgeBaseCreate,
    user: User = Depends(require_user),
    db: AsyncSession = Depends(get_db),
):
    row = UserKnowledgeBase(
        user_id=user.id,
        name=payload.name.strip(),
        description=payload.description.strip(),
        bound_course_id=payload.bound_course_id,
    )
    db.add(row)
    await db.commit()
    await db.refresh(row)
    return _library_payload(row)


@router.patch("/{library_id}")
async def update_library(
    library_id: int,
    payload: KnowledgeBaseUpdate,
    user: User = Depends(require_user),
    db: AsyncSession = Depends(get_db),
):
    row = await _owned_library(db, library_id, user.id)
    if payload.name is not None:
        row.name = payload.name.strip()
    if payload.description is not None:
        row.description = payload.description.strip()
    row.bound_course_id = payload.bound_course_id
    await db.commit()
    await db.refresh(row)
    docs = (
        await db.scalars(
            select(UserKnowledgeDocument).where(
                UserKnowledgeDocument.knowledge_base_id == row.id,
                UserKnowledgeDocument.user_id == user.id,
            )
        )
    ).all()
    return _library_payload(row, list(docs))


@router.delete("/{library_id}")
async def delete_library(
    library_id: int,
    user: User = Depends(require_user),
    db: AsyncSession = Depends(get_db),
):
    row = await _owned_library(db, library_id, user.id)
    source_paths = list(
        await db.scalars(
            select(UserKnowledgeDocument.source_path).where(
                UserKnowledgeDocument.knowledge_base_id == library_id,
                UserKnowledgeDocument.user_id == user.id,
            )
        )
    )
    await db.execute(
        delete(UserKnowledgeChunk).where(
            UserKnowledgeChunk.knowledge_base_id == library_id,
            UserKnowledgeChunk.user_id == user.id,
        )
    )
    await db.execute(
        delete(UserKnowledgeDocument).where(
            UserKnowledgeDocument.knowledge_base_id == library_id,
            UserKnowledgeDocument.user_id == user.id,
        )
    )
    await db.delete(row)
    await db.commit()
    for source_path in source_paths:
        _safe_unlink(source_path or "")
    return {"ok": True}


def _extract_pages(filename: str, raw: bytes) -> list[tuple[int | None, str]]:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(BytesIO(raw))
        return [(index + 1, (page.extract_text() or "").strip()) for index, page in enumerate(reader.pages)]
    if suffix in {".ppt", ".pptx"}:
        if suffix == ".ppt":
            raise HTTPException(status_code=415, detail="旧版 .ppt 请另存为 .pptx 后上传")
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise HTTPException(status_code=503, detail="PPTX 解析组件尚未安装") from exc
        presentation = Presentation(BytesIO(raw))
        return [
            (index + 1, "\n".join(shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()))
            for index, slide in enumerate(presentation.slides)
        ]
    if suffix in {".doc", ".docx"}:
        if suffix == ".doc":
            raise HTTPException(status_code=415, detail="旧版 .doc 请另存为 .docx 后上传")
        try:
            from docx import Document
        except ImportError as exc:
            raise HTTPException(status_code=503, detail="Word 解析组件尚未安装") from exc
        document = Document(BytesIO(raw))
        return [(None, "\n".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()))]
    return [(None, raw.decode("utf-8-sig", errors="replace").strip())]


def _split_pages(pages: list[tuple[int | None, str]], max_chars: int = 1100) -> list[tuple[int | None, str]]:
    chunks: list[tuple[int | None, str]] = []
    for page, text in pages:
        paragraphs = [item.strip() for item in text.replace("\r", "").split("\n") if item.strip()]
        buffer = ""
        for paragraph in paragraphs:
            if buffer and len(buffer) + len(paragraph) + 1 > max_chars:
                chunks.append((page, buffer))
                buffer = ""
            while len(paragraph) > max_chars:
                chunks.append((page, paragraph[:max_chars]))
                paragraph = paragraph[max_chars:]
            buffer = f"{buffer}\n{paragraph}".strip()
        if buffer:
            chunks.append((page, buffer))
    return chunks


async def _persist_document_error(document_id: int, user_id: int, detail: str, *, ocr_status: str | None = None) -> None:
    async with async_session_maker() as db:
        document = await db.scalar(
            select(UserKnowledgeDocument).where(
                UserKnowledgeDocument.id == document_id,
                UserKnowledgeDocument.user_id == user_id,
            )
        )
        if document is None:
            return
        document.status = "error"
        document.error_detail = detail[:1200]
        document.finished_at = datetime.utcnow()
        if ocr_status:
            document.ocr_status = ocr_status
        await db.commit()


async def process_document_task(document_id: int, user_id: int) -> None:
    """进程内队列任务；状态和原文件持久化，进程中断后可由用户安全重试。"""
    try:
        async with async_session_maker() as db:
            document = await db.scalar(
                select(UserKnowledgeDocument).where(
                    UserKnowledgeDocument.id == document_id,
                    UserKnowledgeDocument.user_id == user_id,
                )
            )
            if document is None:
                return
            source = Path(document.source_path).resolve()
            root = _storage_root()
            if not source.is_relative_to(root) or not source.is_file():
                raise RuntimeError("原文件已丢失，无法继续解析")
            document.status = "parsing"
            document.parse_progress = 15
            document.vector_progress = 0
            document.error_detail = ""
            document.started_at = datetime.utcnow()
            document.finished_at = None
            await db.commit()

            raw = source.read_bytes()
            if hashlib.sha256(raw).hexdigest() != document.checksum_sha256:
                raise RuntimeError("原文件校验失败，请删除后重新上传")
            pages = _extract_pages(document.filename, raw)
            chunks = _split_pages(pages)
            if not chunks:
                is_pdf = Path(document.filename).suffix.lower() == ".pdf"
                if is_pdf:
                    raise HTTPException(
                        status_code=422,
                        detail=(
                            "扫描 PDF 未提取到文字；OCR 插件当前未配置。"
                            "可先使用 OCR 工具生成可搜索 PDF，或配置后续 OCR 适配器后安全重试"
                        ),
                    )
                raise HTTPException(status_code=422, detail="没有提取到可检索文字")

            document.page_count = len(pages)
            document.parse_progress = 100
            document.ocr_status = "not_needed"
            await db.execute(
                delete(UserKnowledgeChunk).where(
                    UserKnowledgeChunk.document_id == document.id,
                    UserKnowledgeChunk.user_id == user_id,
                )
            )
            await db.commit()

            vectors: list[list[float] | None] = [None] * len(chunks)
            if has_embedding_key():
                document.status = "vectorizing"
                document.vector_progress = 10
                await db.commit()
                vectors = list(await embed_texts([content for _, content in chunks]))
                # 用户可能在向量化期间删除资料；刷新可阻止后台任务把分片写回已删除资源。
                await db.refresh(document)
                document.vector_progress = 100
                document.status = "ready"
                document.error_detail = ""
            else:
                await db.refresh(document)
                # 关键词索引已就绪，但没有执行语义向量化；保留 0
                # 让界面准确区分两种检索能力。
                document.vector_progress = 0
                document.status = "ready_keyword"
                document.error_detail = "语义向量服务未配置；关键词检索已就绪"

            for (page, content), vector in zip(chunks, vectors):
                db.add(UserKnowledgeChunk(
                    user_id=user_id,
                    knowledge_base_id=document.knowledge_base_id,
                    document_id=document.id,
                    content=content,
                    page=page,
                    embedding=vector,
                ))
            document.finished_at = datetime.utcnow()
            await db.commit()
    except HTTPException as exc:
        is_scanned_pdf = "OCR" in str(exc.detail)
        await _persist_document_error(
            document_id,
            user_id,
            str(exc.detail),
            ocr_status="required_unconfigured" if is_scanned_pdf else None,
        )
    except Exception as exc:
        await _persist_document_error(document_id, user_id, f"资料解析失败：{exc}")


async def mark_interrupted_tasks_failed() -> None:
    """启动恢复：进程中断的任务不伪装为仍在运行，保留原文件供用户重试。"""
    async with async_session_maker() as db:
        rows = (
            await db.scalars(
                select(UserKnowledgeDocument).where(
                    UserKnowledgeDocument.status.in_(("queued", "parsing", "vectorizing"))
                )
            )
        ).all()
        for document in rows:
            document.status = "error"
            document.error_detail = "上次后台任务因服务重启而中断，原文件已保留，可安全重试"
            document.finished_at = datetime.utcnow()
        if rows:
            await db.commit()


@router.post("/{library_id}/documents", status_code=status.HTTP_202_ACCEPTED)
async def upload_document(
    library_id: int,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    user: User = Depends(require_user),
    db: AsyncSession = Depends(get_db),
):
    library = await _owned_library(db, library_id, user.id)
    filename = Path(file.filename or "未命名资料").name
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        raise HTTPException(status_code=415, detail="支持 PDF、PPTX、DOCX、Markdown 与 TXT")
    if suffix in {".ppt", ".doc"}:
        raise HTTPException(status_code=415, detail=f"旧版 {suffix} 请另存为 {'PPTX' if suffix == '.ppt' else 'DOCX'} 后上传")
    raw = await file.read(MAX_FILE_BYTES + 1)
    if len(raw) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail="单个资料不能超过 20MB")
    if not raw:
        raise HTTPException(status_code=400, detail="文件内容为空")

    document = UserKnowledgeDocument(
        user_id=user.id,
        knowledge_base_id=library.id,
        filename=filename,
        media_type=file.content_type or "application/octet-stream",
        size=len(raw),
        status="queued",
        parse_progress=0,
        vector_progress=0,
        ocr_status="pending" if suffix == ".pdf" else "not_needed",
    )
    db.add(document)
    await db.flush()
    source: Path | None = None
    try:
        source = _source_path(user_id=user.id, document_id=document.id, filename=filename)
        source.write_bytes(raw)
        document.source_path = str(source)
        document.checksum_sha256 = hashlib.sha256(raw).hexdigest()
        await db.commit()
        await db.refresh(document)
    except Exception as exc:
        await db.rollback()
        if source is not None:
            _safe_unlink(str(source))
        raise HTTPException(status_code=500, detail=f"资料暂存失败：{exc}") from exc
    background_tasks.add_task(process_document_task, document.id, user.id)
    return _library_payload(library, [document])["documents"][0]


@router.post("/documents/{document_id}/retry", status_code=status.HTTP_202_ACCEPTED)
async def retry_document(
    document_id: int,
    background_tasks: BackgroundTasks,
    user: User = Depends(require_user),
    db: AsyncSession = Depends(get_db),
):
    document = await db.scalar(
        select(UserKnowledgeDocument).where(
            UserKnowledgeDocument.id == document_id,
            UserKnowledgeDocument.user_id == user.id,
        )
    )
    if document is None:
        raise HTTPException(status_code=404, detail="资料不存在")
    if document.status != "error":
        raise HTTPException(status_code=409, detail="只有失败或中断的任务可以重试")
    if document.retry_count >= MAX_SAFE_RETRIES:
        raise HTTPException(status_code=409, detail="已达到安全重试上限，请删除后重新上传原文件")
    source = Path(document.source_path).resolve() if document.source_path else None
    if source is None or not source.is_relative_to(_storage_root()) or not source.is_file():
        raise HTTPException(status_code=409, detail="原文件已丢失，无法重试，请重新上传")
    document.retry_count += 1
    document.status = "queued"
    document.parse_progress = 0
    document.vector_progress = 0
    document.error_detail = ""
    document.started_at = None
    document.finished_at = None
    document.ocr_status = "pending" if Path(document.filename).suffix.lower() == ".pdf" else "not_needed"
    await db.execute(
        delete(UserKnowledgeChunk).where(
            UserKnowledgeChunk.document_id == document.id,
            UserKnowledgeChunk.user_id == user.id,
        )
    )
    await db.commit()
    background_tasks.add_task(process_document_task, document.id, user.id)
    return {"ok": True, "document_id": document.id, "status": "queued", "retry_count": document.retry_count}


@router.delete("/documents/{document_id}")
async def delete_document(
    document_id: int,
    user: User = Depends(require_user),
    db: AsyncSession = Depends(get_db),
):
    document = await db.scalar(
        select(UserKnowledgeDocument).where(
            UserKnowledgeDocument.id == document_id,
            UserKnowledgeDocument.user_id == user.id,
        )
    )
    if document is None:
        raise HTTPException(status_code=404, detail="资料不存在")
    await db.execute(
        delete(UserKnowledgeChunk).where(
            UserKnowledgeChunk.document_id == document_id,
            UserKnowledgeChunk.user_id == user.id,
        )
    )
    await db.delete(document)
    await db.commit()
    _safe_unlink(document.source_path)
    return {"ok": True}


async def search_owned_library(
    db: AsyncSession,
    *,
    user_id: int,
    library_id: int,
    query: str,
    limit: int = 6,
    with_semantic: bool = True,
) -> list[dict]:
    library = await db.scalar(
        select(UserKnowledgeBase).where(
            UserKnowledgeBase.id == library_id,
            UserKnowledgeBase.user_id == user_id,
        )
    )
    if library is None:
        return []
    rows = (
        await db.execute(
            select(UserKnowledgeChunk, UserKnowledgeDocument)
            .join(UserKnowledgeDocument, UserKnowledgeDocument.id == UserKnowledgeChunk.document_id)
            .where(
                UserKnowledgeChunk.knowledge_base_id == library_id,
                UserKnowledgeChunk.user_id == user_id,
                UserKnowledgeDocument.user_id == user_id,
            )
        )
    ).all()
    chunks = [
        Chunk(
            chunk_id=str(chunk.id),
            content=chunk.content,
            source=document.filename,
            page=chunk.page,
            embedding=chunk.embedding,
            meta={"library_id": library.id, "library_name": library.name, "document_id": document.id},
        )
        for chunk, document in rows
    ]
    bm25 = BM25Engine()
    bm25.add(chunks)
    keyword_hits = bm25.search(query, k=limit)
    scores: dict[str, float] = {hit.chunk.chunk_id: hit.score for hit in keyword_hits}
    by_id = {chunk.chunk_id: chunk for chunk in chunks}
    if with_semantic and any(chunk.embedding for chunk in chunks) and has_embedding_key():
        query_vector = await embed_query(query)
        vector = VectorIndex()
        vector.add(chunks)
        for hit in vector.search(query_vector, k=limit):
            scores[hit.chunk.chunk_id] = scores.get(hit.chunk.chunk_id, 0.0) + max(0.0, hit.score) * 2
    ranked = sorted(scores, key=scores.get, reverse=True)[:limit]
    max_score = max((scores[item] for item in ranked), default=1.0)
    return [
        {
            "chunk_id": int(chunk_id),
            "content": by_id[chunk_id].content,
            "source": by_id[chunk_id].source,
            "page": by_id[chunk_id].page,
            "score": scores[chunk_id],
            "relevance_percent": round(scores[chunk_id] / (max_score or 1) * 100),
            "meta": by_id[chunk_id].meta,
        }
        for chunk_id in ranked
    ]


@router.get("/{library_id}/search")
async def search_library(
    library_id: int,
    q: str = Query(min_length=1, max_length=500),
    limit: int = Query(default=6, ge=1, le=12),
    user: User = Depends(require_user),
    db: AsyncSession = Depends(get_db),
):
    await _owned_library(db, library_id, user.id)
    items = await search_owned_library(
        db,
        user_id=user.id,
        library_id=library_id,
        query=q,
        limit=limit,
    )
    return {"query": q, "count": len(items), "items": items}
